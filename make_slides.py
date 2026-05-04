from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT_BLUE = RGBColor(0xD0, 0xE8, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x55, 0x55, 0x55)
ACCENT    = RGBColor(0xE8, 0x6A, 0x10)

W, H = Inches(13.33), Inches(7.5)  # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, l, t, w, h,
             bold=False, italic=False, size=18, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    return txb


def header_bar(slide, title_text, subtitle=None):
    """Dark-blue header bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.25, DARK_BLUE)
    add_text(slide, title_text, 0.3, 0.1, 12.5, 0.8,
             bold=True, size=28, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.85, 12.5, 0.35,
                 size=14, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


def bullet_block(slide, items, l, t, w, h, size=16, color=GRAY, indent="  •  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def footer(slide, text="Osanlou et al., NeurIPS 2024"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, DARK_BLUE)
    add_text(slide, text, 0.3, 7.22, 12.5, 0.25,
             size=10, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 2.8, 13.33, 2.0, MID_BLUE)

add_text(slide, "SWOT-based Simulation of River Discharge",
         0.5, 1.1, 12.3, 1.0, bold=True, size=34, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "with Temporal Graph Neural Networks",
         0.5, 1.95, 12.3, 0.8, bold=True, size=34, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Kevin Osanlou¹  ·  Augusto Getirana²³  ·  Thomas Holmes³  ·  Tristan Cazenave⁴",
         0.5, 2.95, 12.3, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         "¹Talan Research & Innovation  ·  ²SAIC / ³NASA GSFC  ·  ⁴Univ. Paris Dauphine–PSL",
         0.5, 3.45, 12.3, 0.4, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 4.1, 6.3, 0.06, ACCENT)

add_text(slide, "38th Conference on Neural Information Processing Systems (NeurIPS 2024)",
         0.5, 4.3, 12.3, 0.5, size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# Slide 2 – Motivation
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Motivation", "Why estimate river discharge from SWOT?")
footer(slide)

# Left column – SWOT facts
add_rect(slide, 0.4, 1.4, 5.8, 5.4, WHITE)
add_text(slide, "The SWOT Satellite", 0.6, 1.5, 5.4, 0.45,
         bold=True, size=18, color=DARK_BLUE)
bullet_block(slide, [
    "Launched December 2022",
    "Monitors rivers wider than 100 m globally",
    "Measures Water Surface Elevation (WSE), width, and slope",
    "21-day orbital repeat cycle (±78° latitude)",
    "Does NOT measure discharge directly",
], 0.5, 2.0, 5.8, 4.0, size=15, color=GRAY)

# Right column – challenges
add_rect(slide, 6.8, 1.4, 6.1, 5.4, WHITE)
add_text(slide, "Key Challenges", 7.0, 1.5, 5.7, 0.45,
         bold=True, size=18, color=DARK_BLUE)
bullet_block(slide, [
    "Spatial discontinuity – not all reaches observed every day",
    "Temporal discontinuity – sparse 21-day repeat cycle",
    "Discharge is essential for flood risk, water resources management, and hydrological modelling",
    "Reach-scale algorithms are computationally infeasible at global scale",
    "Need: cheap, data-driven regionalization of discharge",
], 6.9, 2.0, 5.8, 4.5, size=15, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# Slide 3 – Problem Statement
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Problem Statement", "Regionalization of river discharge under sparse observations")
footer(slide)

add_rect(slide, 0.4, 1.4, 12.5, 1.5, LIGHT_BLUE)
add_text(slide,
         "Regionalization: inferring hydrological information (discharge) at ungauged or\n"
         "unobserved river reaches from gauged / SWOT-observed nearby reaches.",
         0.6, 1.5, 12.1, 1.2, size=17, color=DARK_BLUE)

add_text(slide, "Formal Setup", 0.5, 3.1, 12.0, 0.4, bold=True, size=18, color=DARK_BLUE)
bullet_block(slide, [
    "Input: temporal sequence of state-of-river graphs over 21 consecutive days",
    "Each node = one SWOT reach (~10 km); edges = connectivity between adjacent reaches",
    "Node features: SWOT observations (sparse), USGS discharge (very sparse, ~3% of nodes), climate variables, static SWORD priors",
    "Target: predict daily discharge at every reach, especially ungauged / SWOT-dark ones",
], 0.5, 3.6, 12.5, 3.0, size=16, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# Slide 4 – Related Work
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Related Work", "Existing approaches to discharge estimation and regionalization")
footer(slide)

cols = [
    ("Drainage-Area Ratio\n(Classic)", [
        "Oldest regionalization method",
        "Ratio of drainage areas between gauged and ungauged sites",
        "Simple but ignores river geometry and temporal dynamics",
    ]),
    ("Data Assimilation", [
        "Combines observations with a physics-based hydrological model",
        "Effective but computationally expensive",
        "Requires complex model construction",
    ]),
    ("ML / LSTM", [
        "Kratzert et al. (2019): rainfall–discharge LSTM outperforms calibrated hydrological models",
        "Ignores spatial (graph) structure of river networks",
    ]),
    ("Graph Neural Networks", [
        "GNNs extend CNNs to non-Euclidean graphs via message passing",
        "Zhao et al., Sun et al.: recurrent GNNs outperform LSTMs for discharge",
        "Struggle with over-smoothing over long distances",
    ]),
]

x_starts = [0.3, 3.55, 6.8, 10.05]
for (title, points), x in zip(cols, x_starts):
    add_rect(slide, x, 1.4, 3.0, 5.6, WHITE)
    add_text(slide, title, x+0.1, 1.5, 2.8, 0.7, bold=True, size=14, color=MID_BLUE)
    bullet_block(slide, points, x+0.05, 2.3, 2.9, 4.3, size=13, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# Slide 5 – Data & Graph Construction
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Data & Graph Construction", "Building temporal sequences of river-state graphs")
footer(slide)

# Node feature table
add_rect(slide, 0.4, 1.4, 5.9, 5.5, WHITE)
add_text(slide, "Node Feature Components", 0.6, 1.5, 5.5, 0.4,
         bold=True, size=17, color=DARK_BLUE)

rows = [
    ("SWOT component", "WSE, width, slope + availability flag (~16% of nodes/day avg.)"),
    ("Discharge component", "USGS gauge value + availability flag (~3% of nodes/day avg.)"),
    ("Climate component", "Temperature, evapotranspiration, cumulative precipitation"),
    ("Static / Prior", "Prior WSE & width, distance to outlet, reach length, channel count (from SWORD)"),
]
y = 2.0
for feat, desc in rows:
    add_rect(slide, 0.5, y, 5.6, 0.05, MID_BLUE)
    add_text(slide, feat, 0.6, y+0.07, 5.3, 0.3, bold=True, size=13, color=DARK_BLUE)
    add_text(slide, desc, 0.6, y+0.38, 5.3, 0.45, size=12, color=GRAY)
    y += 1.0

# Right column – dataset info
add_rect(slide, 6.8, 1.4, 6.1, 5.5, WHITE)
add_text(slide, "Dataset", 7.0, 1.5, 5.7, 0.4, bold=True, size=17, color=DARK_BLUE)
bullet_block(slide, [
    "Foundation: SWOT River Database (SWORD)",
    "6 sub-areas from SWORD basins 73, 74, 78 (30k–60k sq. miles each)",
    "Study period: May 2023 – June 2024",
    "SWOT data: Level-2 single-pass river products",
    "Discharge: USGS WaterService daily API",
    "Climate: Open-Meteo historical weather API",
    "Training signal: 50% of gauged nodes have discharge masked → become labels",
], 6.9, 2.0, 5.8, 4.6, size=14, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# Slide 6 – SWOT-GNN Architecture Overview
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "SWOT-GNN Architecture Overview",
           "Temporal graph neural network pipeline")
footer(slide)

stages = [
    ("1  Input\nFormatting", "Four feature groups (SWOT, discharge, climate, prior) → separate MLPs → concatenated 64-dim graph embedding per day"),
    ("2  ST-Block\n(×2)", "Each ST-Block: (i) node-wise bi-directional LSTM across the 21-day sequence, then (ii) l=3 GraphGPS layers for spatial message passing"),
    ("3  Readout\nMLP", "Per-node MLP maps the final 64-dim embedding to a single discharge prediction for every reach on every day"),
]

x = 0.4
for i, (title, desc) in enumerate(stages):
    add_rect(slide, x, 1.5, 3.8, 4.8, WHITE)
    add_rect(slide, x, 1.5, 3.8, 0.7, MID_BLUE)
    add_text(slide, title, x+0.1, 1.55, 3.6, 0.6,
             bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x+0.15, 2.3, 3.5, 3.7, size=14, color=GRAY)
    if i < 2:
        add_text(slide, "→", x+3.85, 3.6, 0.5, 0.5, bold=True, size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    x += 4.3

add_rect(slide, 0.4, 6.5, 12.5, 0.05, MID_BLUE)
add_text(slide,
         "Final architecture: 2 ST-Blocks × 3 GraphGPS layers + 2-layer bi-LSTM · hidden size 32 · trained on RTX A4000",
         0.4, 6.55, 12.5, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# Slide 7 – ST-Block Detail
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "ST-Block: Temporal + Spatial Processing",
           "The core building block of SWOT-GNN")
footer(slide)

# Temporal side
add_rect(slide, 0.4, 1.4, 5.9, 5.5, WHITE)
add_rect(slide, 0.4, 1.4, 5.9, 0.55, RGBColor(0x2E, 0x6D, 0xA4))
add_text(slide, "Temporal Processing  (bi-LSTM)", 0.55, 1.46, 5.6, 0.45,
         bold=True, size=16, color=WHITE)
bullet_block(slide, [
    "2-layer bidirectional LSTM applied node-wise",
    "Each node receives its own 21-day embedding sequence",
    "Captures short- and long-term temporal dependencies",
    "Dropout (keep=0.5) after the first LSTM layer",
    "Output: temporally-enriched embedding for each node × each day",
], 0.5, 2.05, 5.7, 4.6, size=14, color=GRAY)

# Spatial side
add_rect(slide, 6.8, 1.4, 6.1, 5.5, WHITE)
add_rect(slide, 6.8, 1.4, 6.1, 0.55, ACCENT)
add_text(slide, "Spatial Processing  (GraphGPS ×3)", 6.95, 1.46, 5.8, 0.45,
         bold=True, size=16, color=WHITE)
bullet_block(slide, [
    "3 GraphGPS transformer layers per ST-Block",
    "Each layer = local message passing (GAT) + global attention (FAVOR+)",
    "Local: aggregates features from direct graph neighbours",
    "Global: every node attends to every other node (linear-cost FAVOR+, 4 heads)",
    "Positional encoding: random-walk + graph Laplacian concatenated to each node",
    "Avoids over-smoothing from pure local GNNs",
], 6.9, 2.05, 5.8, 4.6, size=14, color=GRAY)

add_text(slide, "→", 6.3, 3.9, 0.55, 0.5, bold=True, size=28, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# Slide 8 – Input Formatting
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Input Formatting Pipeline",
           "How raw features are embedded before entering ST-Blocks")
footer(slide)

blocks = [
    ("SWOT\n(4 features)", MID_BLUE),
    ("Discharge\n(2 features)", ACCENT),
    ("Climate\n(3 features)", RGBColor(0x2A, 0x9D, 0x5C)),
    ("Prior\n(6 features)", RGBColor(0x8B, 0x44, 0xAC)),
]

x = 0.5
for label, col in blocks:
    add_rect(slide, x, 1.5, 2.7, 0.8, col)
    add_text(slide, label, x+0.05, 1.55, 2.6, 0.7,
             bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow down
    add_text(slide, "↓", x+1.1, 2.35, 0.5, 0.4, bold=True, size=20, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    # LayerNorm box
    add_rect(slide, x, 2.75, 2.7, 0.5, LIGHT_BLUE)
    add_text(slide, "LayerNorm", x+0.05, 2.8, 2.6, 0.4, size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "↓", x+1.1, 3.3, 0.5, 0.4, bold=True, size=20, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    # MLP box
    add_rect(slide, x, 3.7, 2.7, 0.5, col)
    add_text(slide, "MLP → 16-dim", x+0.05, 3.75, 2.6, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    x += 3.1

add_text(slide, "CONCATENATE  →  64-dim  →  LayerNorm  →  Graph Embedding (64-dim per node)",
         1.0, 4.4, 11.3, 0.5, bold=True, size=16, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.0, 4.35, 11.3, 0.05, DARK_BLUE)
add_rect(slide, 1.0, 4.9, 11.3, 0.05, DARK_BLUE)

add_text(slide,
         "Each of the 4 feature groups is independently normalized and projected to 16 dimensions.\n"
         "Their concatenation (64-dim) is the node embedding fed into the ST-Blocks.",
         0.5, 5.1, 12.3, 1.9, size=14, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# Slide 9 – Experiments & Results
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF8, 0xFD))
header_bar(slide, "Experiments & Results",
           "KGE evaluation on USGS gauges across two test basins")
footer(slide)

# Metric box
add_rect(slide, 0.4, 1.4, 4.5, 2.0, WHITE)
add_text(slide, "Evaluation Metric", 0.6, 1.5, 4.1, 0.4, bold=True, size=16, color=DARK_BLUE)
add_text(slide,
         "Kling–Gupta Efficiency (KGE)\n∈ (−∞, 1]  (higher is better)\n"
         "Combines correlation, bias, and variability ratio",
         0.6, 1.95, 4.1, 1.3, size=13, color=GRAY)

# Baselines
add_rect(slide, 5.2, 1.4, 7.7, 2.0, WHITE)
add_text(slide, "Compared Models", 5.4, 1.5, 7.3, 0.4, bold=True, size=16, color=DARK_BLUE)
bullet_block(slide, [
    "SWOT-GNN  (proposed)  — temporal GNN",
    "GPS-GNN — 6 GraphGPS layers, no temporal component",
    "LSTM — 4-layer bi-directional LSTM, no spatial component",
    "Drainage-area ratio — classic baseline using non-masked gauges",
], 5.3, 1.95, 7.5, 1.3, size=13, color=GRAY)

# Results table
add_rect(slide, 0.4, 3.55, 12.5, 0.5, MID_BLUE)
for txt, x in [("Model", 0.55), ("Basin 1 – Min KGE", 2.5), ("Basin 1 – Mean KGE", 5.2),
               ("Basin 2 – Mean KGE", 7.9), ("Overall", 10.6)]:
    add_text(slide, txt, x, 3.6, 2.5, 0.4, bold=True, size=13, color=WHITE)

rows_data = [
    ("SWOT-GNN", "0.10", "~0.32", "~0.28", "Best"),
    ("GPS-GNN",  "0.05", "~0.25", "~0.22", "2nd"),
    ("LSTM",     "<0",   "~0.15", "~0.10", "3rd"),
    ("DAR",      "<0",   "~0.10", "~0.05", "Weakest"),
]
row_colors = [LIGHT_BLUE, WHITE, LIGHT_BLUE, WHITE]
y = 4.1
for (model, mn, m1, m2, rank), bg in zip(rows_data, row_colors):
    add_rect(slide, 0.4, y, 12.5, 0.55, bg)
    for txt, x in [(model, 0.55), (mn, 2.5), (m1, 5.2), (m2, 7.9), (rank, 10.6)]:
        add_text(slide, txt, x, y+0.07, 2.4, 0.4, size=13, color=GRAY)
    y += 0.56

add_text(slide,
         "SWOT-GNN outperforms all baselines. Low-coverage gauges remain challenging for all models.",
         0.4, 7.0, 12.5, 0.4, size=13, color=DARK_BLUE, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# Slide 10 – Conclusions & Future Work
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 1.1, 13.33, 5.3, RGBColor(0xF4, 0xF8, 0xFD))

add_text(slide, "Conclusions & Future Work", 0.5, 0.15, 12.3, 0.85,
         bold=True, size=30, color=WHITE)

add_rect(slide, 0.5, 1.3, 5.9, 4.8, WHITE)
add_rect(slide, 0.5, 1.3, 5.9, 0.55, MID_BLUE)
add_text(slide, "Key Contributions", 0.65, 1.35, 5.6, 0.45, bold=True, size=16, color=WHITE)
bullet_block(slide, [
    "SWOT-GNN: first temporal GNN tailored to SWOT data for basin-scale discharge simulation",
    "End-to-end pipeline: SWORD + SWOT + USGS + Open-Meteo → temporal graph sequences",
    "Handles spatial & temporal data discontinuity via graph regionalization",
    "Outperforms drainage-area ratio, GPS-GNN (no temporal), and LSTM (no spatial) on KGE",
    "Computationally cheap vs. physics-based hydrological models",
], 0.6, 1.95, 5.7, 4.0, size=14, color=GRAY)

add_rect(slide, 6.9, 1.3, 6.0, 4.8, WHITE)
add_rect(slide, 6.9, 1.3, 6.0, 0.55, ACCENT)
add_text(slide, "Future Directions", 7.05, 1.35, 5.7, 0.45, bold=True, size=16, color=WHITE)
bullet_block(slide, [
    "Incorporate additional hydrological parameters (soil type, elevation, land use) to improve performance in SWOT-sparse areas",
    "Integrate upcoming SWOT reach-scale discharge products for model fine-tuning",
    "Extend to global river networks beyond the CONUS test basins",
    "Explore multi-task learning (WSE + discharge jointly)",
    "Investigate uncertainty quantification for operational flood forecasting",
], 7.0, 1.95, 5.8, 4.0, size=14, color=GRAY)

add_rect(slide, 0, 6.4, 13.33, 1.1, DARK_BLUE)
add_text(slide,
         "SWOT-GNN provides a scalable, data-driven path to global river discharge estimation\n"
         "by fusing satellite observations, in-situ gauges, and climate data via temporal graph learning.",
         0.5, 6.45, 12.3, 0.9, size=14, color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────
out = r"c:\Disk_D\My_drive\Code\Project_code_2025_2026\smart_hs_mekong_mega_reservoirs\SWOT_GNN_overview.pptx"
prs.save(out)
print(f"Saved: {out}")
